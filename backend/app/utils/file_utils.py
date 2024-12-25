import os
import io
from google.auth.transport.requests import Request
from google.oauth2.credentials import Credentials
from google_auth_oauthlib.flow import InstalledAppFlow
from googleapiclient.discovery import build
from googleapiclient.http import MediaIoBaseDownload
from googleapiclient.errors import HttpError
from docx import Document
import pandas as pd
from pptx import Presentation
from langchain_community.document_loaders import PyPDFLoader

SCOPES = ["https://www.googleapis.com/auth/drive.readonly"]

def getCredentials():
    creds = None
    # The file token.json stores the user's access and refresh tokens, and is
    # created automatically when the authorization flow completes for the first
    # time.
    if os.path.exists("token.json"):
        creds = Credentials.from_authorized_user_file("token.json", SCOPES)
    # If there are no (valid) credentials available, let the user log in.
    if not creds or not creds.valid:
        if creds and creds.expired and creds.refresh_token:
            creds.refresh(Request())
        else:
            flow = InstalledAppFlow.from_client_secrets_file(
                "credentials.json", SCOPES
            )
            creds = flow.run_local_server(port=0)
        # Save the credentials for the next run
        with open("token.json", "w") as token:
            token.write(creds.to_json())
    return creds

def getFilesInFolder(folder_id):
    """Fetches files from the specified Google Drive folder."""
    files = []
    try:
        service = build("drive", "v3", credentials=getCredentials())
        page_token = None
        while True:
            # Call the Drive v3 API
            results = (
                service.files()
                .list(q=f"'{folder_id}' in parents",  # Specify folder_id
                      spaces="drive",
                      fields="nextPageToken, files(id, name)",
                      pageToken=page_token)
                .execute()
            )
            items = results.get("files", [])
            for item in items:
                files.append(item)
            page_token = results.get("nextPageToken", None)
            if not page_token:
                break
    except HttpError as error:
        print(f"An error occurred: {error}")
    return files

def download_file(service, file_id, file_name):
    """Download a file from Google Drive."""
    request = service.files().get_media(fileId=file_id)
    fh = io.FileIO(file_name, 'wb')
    downloader = MediaIoBaseDownload(fh, request)
    done = False
    while done is False:
        status, done = downloader.next_chunk()
        print(f"Download {int(status.progress() * 100)}%.")
    return file_name

def load_hidden_documents_from_drive(files):
    all_texts = []
    for file in files:
        file_name = file['name']
        file_id = file['id']
        
        # Download the file from Google Drive
        print(f"Downloading file: {file_name}")
        downloaded_file_name = download_file(service, file_id, file_name)
        
        # Read the downloaded file content and extract text
        if downloaded_file_name.endswith(".pdf"):
            loader = PyPDFLoader(downloaded_file_name)
            all_texts.extend([page.page_content for page in loader.load_and_split()])
        elif downloaded_file_name.endswith(".docx"):
            doc = Document(downloaded_file_name)
            all_texts.append("\n".join(p.text for p in doc.paragraphs))
        elif downloaded_file_name.endswith(".txt"):
            with open(downloaded_file_name, "r", encoding="utf-8") as f:
                all_texts.append(f.read())
        elif downloaded_file_name.endswith((".xlsx", ".xls")):
            all_texts.append(pd.read_excel(downloaded_file_name).to_string(index=False))
        elif downloaded_file_name.endswith(".pptx"):
            presentation = Presentation(downloaded_file_name)
            all_texts.append("\n".join(shape.text for slide in presentation.slides for shape in slide.shapes if shape.has_text_frame))
        
        # Delete the file after processing
        os.remove(downloaded_file_name)
    
    return all_texts

if __name__ == "__main__":
    # Specify your folder ID here
    folder_id = os.getenv("folder_id")
    
    # Fetch the files from the specified folder
    service = build("drive", "v3", credentials=getCredentials())
    files = getFilesInFolder(folder_id)
    
    # Process the files to extract text and build all_texts
    all_texts = load_hidden_documents_from_drive(files)
    print("Collected all texts from files:")
    print(all_texts)

