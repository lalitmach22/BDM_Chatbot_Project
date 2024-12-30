import os, re, mimetypes
from docx import Document
import pandas as pd
from pptx import Presentation
from langchain_community.document_loaders import PyPDFLoader

def clean_text(text):
    """Clean text by removing extra spaces, fixing punctuation, and standardizing newlines."""
    text = re.sub(r'\s+', ' ', text)  # Replace multiple spaces with a single space
    text = re.sub(r'(?<=[a-zA-Z])(?=[A-Z])', ' ', text)  # Add space between words when stuck together
    text = re.sub(r'(?<=[.,?!;])(?=[a-zA-Z])', ' ', text)  # Add space after punctuation if missing
    text = re.sub(r'\.\s+', '.\n', text)  # Add newlines after sentences
    text = re.sub(r'(?<=:)\s+', '\n', text)  # Add newlines after colons
    text = text.strip()  # Remove leading/trailing spaces
    return text

def load_pdf(file_path):
    """Load PDF and return text from each page."""
    try:
        loader = PyPDFLoader(file_path)
        all_texts = []
        pages = loader.load_and_split()
        
        for page in pages:
            all_texts.append(page.page_content)
        return all_texts
    except Exception as e:
        print(f"Error loading PDF {file_path}: {e}")
        return []

def load_word(file_path, chunk_size=10):
    """Load Word document and return text in chunks."""
    try:
        doc = Document(file_path)
        all_texts = []
        paragraphs = [paragraph.text for paragraph in doc.paragraphs]
        for i in range(0, len(paragraphs), chunk_size):
            chunk = "\n".join(paragraphs[i:i+chunk_size])
            all_texts.append(chunk)
        return all_texts
    except Exception as e:
        print(f"Error loading Word document {file_path}: {e}")
        return []

def load_text(file_path, chunk_size=1024):
    """Load plain text file and return text in chunks."""
    try:
        with open(file_path, "r", encoding="utf-8") as file:
            all_texts = []
            while chunk := file.read(chunk_size):
                all_texts.append(chunk)
        return all_texts
    except Exception as e:
        print(f"Error loading text file {file_path}: {e}")
        return []

def load_excel(file_path, chunk_size=1000):
    """Load Excel file and return text in chunks."""
    try:
        excel_data = pd.read_excel(file_path, chunksize=chunk_size)
        all_texts = []
        for chunk in excel_data:
            all_texts.append(chunk.to_string(index=False))
        return all_texts
    except Exception as e:
        print(f"Error loading Excel file {file_path}: {e}")
        return []

def load_csv(file_path, chunk_size=1000):
    """Load CSV file and return text in chunks."""
    try:
        csv_data = pd.read_csv(file_path, chunksize=chunk_size)
        all_texts = []
        for chunk in csv_data:
            all_texts.append(chunk.to_string(index=False))
        return all_texts
    except Exception as e:
        print(f"Error loading CSV file {file_path}: {e}")
        return []

def load_pptx(file_path):
    """Load PowerPoint file and return text from each slide."""
    try:
        presentation = Presentation(file_path)
        all_texts = []
        for slide in presentation.slides:
            slide_text = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    slide_text.append(shape.text)
            all_texts.append("\n".join(slide_text))
        return all_texts
    except Exception as e:
        print(f"Error loading PowerPoint file {file_path}: {e}")
        return []

def load_hidden_documents(directory="hidden_docs"):
    """Load all supported files from a directory and return their content in chunks."""
    all_texts = []
    for filename in os.listdir(directory):
        file_path = os.path.join(directory, filename)
        mime_type, _ = mimetypes.guess_type(file_path)

        # Handle different file types based on extension
        try:
            if filename.endswith(".pdf"):
                all_texts.extend(load_pdf(file_path))
            elif filename.endswith(".docx"):
                all_texts.extend(load_word(file_path))
            elif filename.endswith(".txt"):
                all_texts.extend(load_text(file_path))
            elif filename.endswith(('.xlsx', '.xls')):
                all_texts.extend(load_excel(file_path))
            elif filename.endswith(".csv"):
                all_texts.extend(load_csv(file_path))
            elif filename.endswith(".pptx"):
                all_texts.extend(load_pptx(file_path))
        except Exception as e:
            print(f"Error processing file {filename}: {e}")

    # Clean the collected texts
    cleaned_texts = [clean_text(text) for text in all_texts]
    return cleaned_texts
