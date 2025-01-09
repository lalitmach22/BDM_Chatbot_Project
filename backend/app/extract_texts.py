import os
import re
import pandas as pd
from docx import Document
from pptx import Presentation
import mimetypes
from zipfile import ZipFile
from langchain_community.document_loaders import PyPDFLoader
import logging
import shutil

# Configure logging
logging.basicConfig(level=logging.INFO, format='%(asctime)s - %(levelname)s - %(message)s')
logger = logging.getLogger(__name__)

def clean_text(text):
    
    """
    Clean text by replacing multiple spaces with a single space, fixing words broken by line breaks, 
    fixing punctuation followed by words with no space, standardizing newlines, and performing additional cleanup.

    Parameters
    ----------
    text : str
        Input text to be cleaned.

    Returns
    -------
    str
        Cleaned text.
    """
    text = re.sub(r'\s+', ' ', text)
    # Fix words broken by line breaks or formatting
    text = re.sub(r'(?<=[a-zA-Z])(?=[A-Z])', ' ', text)  # Add space between words when they are stuck together    
    # Fix punctuation followed by words with no space
    text = re.sub(r'(?<=[.,?!;])(?=[a-zA-Z])', ' ', text)  # Add space after punctuation if missing
    # Standardize newlines for better formatting
    text = re.sub(r'\.\s+', '.\n', text)  # Add newlines after sentences
    text = re.sub(r'(?<=:)\s+', '\n', text)  # Add newlines after colons
    # Additional cleanup (if needed)
    text = text.strip()  # Remove leading and trailing whitespace
    return text

def load_hidden_documents(directory, files =None):
    """Load all supported file types from a directory and return their content."""
    all_texts = []

    # If files are provided, filter the directory for those files
    files_to_process = files if files else os.listdir(directory)

    for filename in os.listdir(directory):
        file_path = os.path.join(directory, filename)
        mime_type, _ = mimetypes.guess_type(file_path)

        try:
            # Handle PDF files
            if filename.endswith(".pdf"):
                loader = PyPDFLoader(file_path)
                pages = loader.load_and_split()
                all_texts.extend([page.page_content for page in pages])
                logger.info(f"Processed PDF file: {filename}")

            # Handle Word files (.docx)
            elif filename.endswith(".docx"):
                doc = Document(file_path)
                text = "\n".join([paragraph.text for paragraph in doc.paragraphs])
                all_texts.append(text)
                logger.info(f"Processed DOCX file: {filename}")

            # Handle Text files (.txt)
            elif filename.endswith(".txt"):
                with open(file_path, "r", encoding="utf-8") as file:
                    all_texts.append(file.read())
                logger.info(f"Processed TXT file: {filename}")

            # Handle CSV files (.csv)
            elif filename.endswith(".csv"):
                csv_data = pd.read_csv(file_path)
                text = csv_data.to_string(index=False)
                all_texts.append(text)
                logger.info(f"Processed CSV file: {filename}")

            # Handle PowerPoint files (.pptx)
            elif filename.endswith(".pptx"):
                presentation = Presentation(file_path)
                for slide in presentation.slides:
                    slide_text = []
                    for shape in slide.shapes:
                        if shape.has_text_frame:
                            slide_text.append(shape.text)
                    all_texts.append("\n".join(slide_text))
                logger.info(f"Processed PPTX file: {filename}")

            # Handle ZIP files (.zip)
            elif filename.endswith(".zip"):
                temp_dir = "temp_extracted"
                try:
                    with ZipFile(file_path, 'r') as zip_ref:
                        zip_ref.extractall(temp_dir)
                        all_texts.extend(load_hidden_documents(temp_dir))
                finally:
                    shutil.rmtree(temp_dir)  # Clean up temporary directory
                logger.info(f"Processed ZIP file: {filename}")

            # Handle unknown file types (fallback to text-based reading)
            elif mime_type and mime_type.startswith("text"):
                with open(file_path, "r", encoding="utf-8") as file:
                    all_texts.append(file.read())
                logger.info(f"Processed text file: {filename}")

        except Exception as e:
            logger.error(f"Failed to process {filename}: {e}")
    cleaned_texts = [clean_text(text) for text in all_texts]
    return cleaned_texts